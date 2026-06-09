"""统一文档解析器 — 支持 PDF/DOCX/XLSX/PPTX/TXT/MD/CSV/JSON/HTML。"""

import csv
import io
import json
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class ParsedDocument:
    """解析结果：按页/段落组织的文本块列表。"""

    def __init__(self, filename: str, pages: list[dict]):
        """
        Args:
            filename: 原始文件名
            pages: [{"page": int, "text": str}, ...] 每页/每段的文本
        """
        self.filename = filename
        self.pages = pages

    @property
    def full_text(self) -> str:
        return "\n\n".join(p["text"] for p in self.pages)


def parse_pdf(file_path: str) -> ParsedDocument:
    """逐页提取 PDF 文本。"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append({"page": i + 1, "text": text.strip()})
    doc.close()
    return ParsedDocument(Path(file_path).name, pages)


def parse_docx(file_path: str) -> ParsedDocument:
    """段落级提取 Word 文本。"""
    from docx import Document

    doc = Document(file_path)
    # 按自然段落分组，每 5 个段落合并为一页（避免 chunk 过碎）
    all_paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages = []
    chunk_size = 5
    for i in range(0, len(all_paragraphs), chunk_size):
        group = all_paragraphs[i: i + chunk_size]
        pages.append({"page": i // chunk_size + 1, "text": "\n".join(group)})
    return ParsedDocument(Path(file_path).name, pages)


def parse_xlsx(file_path: str) -> ParsedDocument:
    """逐行提取 Excel 文本。"""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    pages = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            row_str = " | ".join(str(cell) for cell in row if cell is not None)
            if row_str.strip():
                rows_text.append(row_str)
        if rows_text:
            pages.append({"page": 1, "text": f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text)})
    wb.close()
    if not pages:
        pages = [{"page": 1, "text": "(空文件)"}]
    return ParsedDocument(Path(file_path).name, pages)


def parse_pptx(file_path: str) -> ParsedDocument:
    """逐页 slide 提取文本。"""
    from pptx import Presentation

    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            pages.append({"page": i + 1, "text": "\n".join(texts)})
    return ParsedDocument(Path(file_path).name, pages)


def parse_txt(file_path: str) -> ParsedDocument:
    """直接读取纯文本文件。"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return ParsedDocument(Path(file_path).name, [{"page": 1, "text": text.strip()}])


def parse_csv(file_path: str) -> ParsedDocument:
    """CSV 逐行读取。"""
    rows = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            line = " | ".join(cell for cell in row if cell.strip())
            if line.strip():
                rows.append(line)
    text = "\n".join(rows) if rows else "(空文件)"
    return ParsedDocument(Path(file_path).name, [{"page": 1, "text": text}])


def parse_json(file_path: str) -> ParsedDocument:
    """JSON 格式化后读取。"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    text = json.dumps(data, ensure_ascii=False, indent=2)
    return ParsedDocument(Path(file_path).name, [{"page": 1, "text": text}])


def parse_html(file_path: str) -> ParsedDocument:
    """BeautifulSoup 提取纯文本。"""
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "lxml")
    # 移除 script 和 style
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    return ParsedDocument(Path(file_path).name, [{"page": 1, "text": text}])


# 解析器映射表
PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
    ".txt": parse_txt,
    ".md": parse_txt,
    ".csv": parse_csv,
    ".json": parse_json,
    ".html": parse_html,
}


def parse_document(file_path: str) -> ParsedDocument:
    """根据文件类型分发到对应的解析器。

    Args:
        file_path: 文件的本地路径

    Returns:
        ParsedDocument 对象

    Raises:
        ValueError: 不支持的文件类型
        RuntimeError: 解析失败
    """
    suffix = Path(file_path).suffix.lower()
    parser = PARSERS.get(suffix)
    if parser is None:
        raise ValueError(f"不支持的文件类型: {suffix}")

    try:
        result = parser(file_path)
        logger.info(f"解析成功: {Path(file_path).name}, 共 {len(result.pages)} 页/段")
        return result
    except Exception as e:
        logger.error(f"解析失败: {Path(file_path).name}, 错误: {e}")
        raise RuntimeError(f"文档解析失败: {e}")
